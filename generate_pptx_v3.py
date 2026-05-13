"""Generate financial model knowledge graph demo presentation V3."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Color Constants ---
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
GOLD = RGBColor(0xD4, 0xA8, 0x43)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GOLD_LIGHT = RGBColor(0xFD, 0xF0, 0xD0)
NAVY_LIGHT = RGBColor(0x2C, 0x3E, 0x6B)
GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)

FONT_CN = "Microsoft YaHei"
FONT_NUM = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 13

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ======================== Helper Functions ========================

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = line_width or Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_rect_straight(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = line_width or Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size, color, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(tf, text, font_size, color, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN, space_before=None, space_after=None):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return p


def add_slide_number(slide, num):
    add_textbox(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                f"{num}/{TOTAL_SLIDES}", 14, RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, fill_color, border_color, title, body, title_size=22, body_size=16):
    """Add a card with title and body text."""
    shape = add_rect(slide, left, top, width, height, fill_color, border_color, Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(title_size)
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = FONT_CN
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.paragraphs[0].space_after = Pt(8)
    add_paragraph(tf, body, body_size, DARK_TEXT, font_name=FONT_CN, space_before=4)
    return shape


def add_icon_circle(slide, left, top, size, text, bg_color=NAVY, text_color=WHITE):
    """Add a circle with number inside."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = text_color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = FONT_NUM
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape


def add_gold_callout(slide, left, top, width, height, text, font_size=18):
    shape = add_rect(slide, left, top, width, height, GOLD_LIGHT, GOLD, Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.bold = False
    tf.paragraphs[0].font.name = FONT_CN
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    return shape


# ======================== Slide 1: Title ========================
def build_slide_1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Gold horizontal line above title
    add_rect_straight(slide, Inches(4.0), Inches(2.2), Inches(5.333), Pt(3), GOLD)

    # Title
    add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10.333), Inches(1.2),
                "财务模型知识图谱系统", 44, GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    # Gold horizontal line below title
    add_rect_straight(slide, Inches(4.0), Inches(3.7), Inches(5.333), Pt(3), GOLD)

    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10.333), Inches(1.0),
                '从"文件"到"资产" — 给财务模型装上 X 光片', 24, WHITE, alignment=PP_ALIGN.CENTER)

    # Tagline
    add_textbox(slide, Inches(1.5), Inches(6.0), Inches(10.333), Inches(0.6),
                "可理解 · 可操作 · 可追溯", 20, GOLD, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 1)


# ======================== Slide 2: Anti-intuitive Question ========================
def build_slide_2():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Heading
    add_textbox(slide, Inches(1.0), Inches(0.5), Inches(11.333), Inches(0.8),
                "做财务模型，最值钱的资产是什么？", 36, NAVY, bold=True)

    # Gold accent line
    add_rect_straight(slide, Inches(1.0), Inches(1.35), Inches(3.0), Pt(3), GOLD)

    # Three bullet points
    bullets = [
        "不是 Excel 文件 —— 文件会散",
        "不是模板 —— 模板会旧",
        "是团队脑子里的逻辑 —— 为什么这样建模、参数怎么取、依赖关系是什么",
    ]
    y = Inches(1.8)
    for i, b in enumerate(bullets):
        prefix = f"{i+1}."
        add_textbox(slide, Inches(1.2), y, Inches(0.6), Inches(0.5),
                    prefix, 22, GOLD, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_NUM)
        add_textbox(slide, Inches(1.8), y, Inches(10.0), Inches(0.5),
                    b, 22, DARK_TEXT)
        y += Inches(0.7)

    # Callout box
    add_gold_callout(slide, Inches(1.0), Inches(4.3), Inches(11.333), Inches(0.8),
                     "“”但这些知识，项目一结束，就散了。”", 18)

    # Bottom text
    add_textbox(slide, Inches(8.0), Inches(6.6), Inches(4.333), Inches(0.5),
                "「今天演示的东西，解决的就是这个问题。」",
                16, GRAY, alignment=PP_ALIGN.RIGHT)

    add_slide_number(slide, 2)


# ======================== Slide 3: Big Picture ========================
def build_slide_3():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1.0), Inches(0.5), Inches(11.333), Inches(0.8),
                "从文件到资产", 36, NAVY, bold=True)
    add_rect_straight(slide, Inches(1.0), Inches(1.35), Inches(2.0), Pt(3), GOLD)

    # Three-stage flow
    # Box 1: 客户给的 Excel
    b1 = add_rect(slide, Inches(1.0), Inches(2.0), Inches(3.0), Inches(1.2), NAVY, GOLD, Pt(2))
    tf = b1.text_frame
    tf.paragraphs[0].text = "客户给的 Excel"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = FONT_CN
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Arrow 1
    add_textbox(slide, Inches(4.1), Inches(2.3), Inches(0.8), Inches(0.6),
                "→", 36, GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    # Box 2: 知识图谱
    b2 = add_rect(slide, Inches(5.0), Inches(1.8), Inches(3.333), Inches(1.6), GOLD_LIGHT, GOLD, Pt(3))
    tf = b2.text_frame
    tf.paragraphs[0].text = "知识图谱"
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = FONT_CN
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_paragraph(tf, "看得见 · 改得起 · 留得住", 16, GOLD, bold=True,
                  alignment=PP_ALIGN.CENTER, space_before=8)

    # Arrow 2
    add_textbox(slide, Inches(8.4), Inches(2.3), Inches(0.8), Inches(0.6),
                "→", 36, GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    # Box 3: 可交付的 Excel
    b3 = add_rect(slide, Inches(9.3), Inches(2.0), Inches(3.0), Inches(1.2), NAVY, GOLD, Pt(2))
    tf = b3.text_frame
    tf.paragraphs[0].text = "可交付的 Excel"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = FONT_CN
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Three sentences
    sentences = [
        "输入：任何 .xlsx 文件，不需要模板化",
        "中间：自动提取公式依赖、指标关系、逻辑链",
        "输出：改完导回 Excel，还是一个能用的财务模型",
    ]
    y = Inches(4.0)
    for s in sentences:
        add_textbox(slide, Inches(1.5), y, Inches(10.333), Inches(0.4),
                    s, 18, DARK_TEXT)
        y += Inches(0.45)

    # Key message box
    add_gold_callout(slide, Inches(1.0), Inches(5.6), Inches(11.333), Inches(1.2),
                     "“不是替代 Excel — 咨询顾问的交付物还是 Excel。\n这是中间的理解层和操作层。”", 18)

    add_slide_number(slide, 3)


# ======================== Slide 4: Core Concepts ========================
def build_slide_4():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1.0), Inches(0.5), Inches(11.333), Inches(0.8),
                "知识图谱的四层结构", 36, NAVY, bold=True)
    add_rect_straight(slide, Inches(1.0), Inches(1.35), Inches(2.0), Pt(3), GOLD)

    # Left column header
    add_rect_straight(slide, Inches(1.0), Inches(1.7), Inches(5.5), Inches(0.55), NAVY)
    add_textbox(slide, Inches(1.1), Inches(1.72), Inches(5.3), Inches(0.5),
                "物理层 — Excel 的原始要素", 20, WHITE, bold=True)

    # Left column content
    left_items = [
        ("Sheet", "Excel 工作表 — 模型的“容器”"),
        ("Cell", "单元格 — 数据的最小单位，带公式和值"),
    ]
    y = Inches(2.5)
    for label, desc in left_items:
        add_textbox(slide, Inches(1.2), y, Inches(1.2), Inches(0.5),
                    label, 22, GOLD, bold=True, font_name=FONT_NUM)
        add_textbox(slide, Inches(2.4), y, Inches(3.8), Inches(0.5),
                    desc, 20, DARK_TEXT)
        y += Inches(0.7)

    # Right column header
    add_rect_straight(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(0.55), GOLD)
    add_textbox(slide, Inches(6.9), Inches(1.72), Inches(5.3), Inches(0.5),
                "业务层 — 从 Cell 抽象出来的", 20, NAVY, bold=True)

    # Right column content
    right_items = [
        ("Table", "逻辑表 — 同一业务主题的 Cell 集合（如“资金筹措及还本付息表”）"),
        ("Indicator", "指标 — 有业务意义的行（如“动态总投资”），是业务语言，不是单元格坐标"),
    ]
    y = Inches(2.5)
    for label, desc in right_items:
        add_textbox(slide, Inches(7.0), y, Inches(1.4), Inches(0.5),
                    label, 22, NAVY, bold=True, font_name=FONT_NUM)
        add_textbox(slide, Inches(8.4), y, Inches(3.8), Inches(0.5),
                    desc, 20, DARK_TEXT)
        y += Inches(0.7)

    # Bottom insight - italic box
    shape = add_rect(slide, Inches(1.0), Inches(4.5), Inches(11.333), Inches(1.8),
                     LIGHT_GRAY, GOLD, Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "“Sheet 和 Cell 是物理的，Table 和 Indicator 是虚拟的。"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.name = FONT_CN
    add_paragraph(tf, "虚拟层与业务绑定 — 这就是知识图谱的核心价值。”",
                  18, NAVY, font_name=FONT_CN, bold=False, space_before=6)

    add_slide_number(slide, 4)


# ======================== Slide 5: Four Scenarios ========================
def build_slide_5():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1.0), Inches(0.5), Inches(11.333), Inches(0.8),
                "我们怎么帮到你？", 36, NAVY, bold=True)
    add_rect_straight(slide, Inches(1.0), Inches(1.35), Inches(2.0), Pt(3), GOLD)

    # 2x2 grid of cards
    scenarios = [
        ("1", "新项目组接手", "3天理解变成3小时"),
        ("2", "方案汇报", "让客户“看得见”你的专业"),
        ("3", "参数调整", "改了敢交"),
        ("4", "客户问答", "从翻表到秒答"),
    ]
    positions = [
        (Inches(1.0), Inches(1.8)),
        (Inches(7.0), Inches(1.8)),
        (Inches(1.0), Inches(4.2)),
        (Inches(7.0), Inches(4.2)),
    ]
    card_w, card_h = Inches(5.3), Inches(2.0)

    for i, (num, title, val) in enumerate(scenarios):
        left, top = positions[i]
        # Card background
        add_rect(slide, left, top, card_w, card_h, WHITE, GOLD, Pt(2))
        # Number circle
        add_icon_circle(slide, left + Inches(0.3), top + Inches(0.4), Inches(0.6), num)
        # Title
        add_textbox(slide, left + Inches(1.1), top + Inches(0.35), Inches(3.8), Inches(0.6),
                    title, 24, NAVY, bold=True)
        # Value line
        add_textbox(slide, left + Inches(1.1), top + Inches(1.0), Inches(3.8), Inches(0.5),
                    val, 20, GOLD, bold=True)

    # Bottom text
    add_textbox(slide, Inches(1.0), Inches(6.5), Inches(11.333), Inches(0.4),
                "“每个场景都是真实的咨询工作流”",
                16, GRAY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 5)


# ======================== Slide 6: Scenario 1 ========================
def build_slide_6():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1.0), Inches(0.5), Inches(11.333), Inches(0.8),
                "场景 1：新项目组接手 — 3天变成3小时", 32, NAVY, bold=True)
    add_rect_straight(slide, Inches(1.0), Inches(1.2), Inches(2.0), Pt(3), GOLD)

    # Left column - 业务场景
    add_rect_straight(slide, Inches(1.0), Inches(1.6), Inches(5.5), Inches(0.5), NAVY)
    add_textbox(slide, Inches(1.1), Inches(1.62), Inches(5.3), Inches(0.45),
                "业务场景", 18, WHITE, bold=True)

    add_textbox(slide, Inches(1.2), Inches(2.3), Inches(5.2), Inches(2.0),
                "新项目启动，拿到客户 50+ Sheet 的财务模型。\n初级顾问需要多久理解？",
                18, DARK_TEXT)

    # Right column - 演示要点
    add_rect_straight(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.5), GOLD)
    add_textbox(slide, Inches(6.9), Inches(1.62), Inches(5.3), Inches(0.45),
                "演示要点", 18, NAVY, bold=True)

    points = [
        "上传文件 → 自动解析 → 图谱展开",
        "点击关键指标，一键追溯上下游",
        "49张表、2968个指标一目了然",
    ]
    y = Inches(2.3)
    for p in points:
        add_textbox(slide, Inches(7.0), y, Inches(5.0), Inches(0.5),
                    f"•  {p}", 18, DARK_TEXT)
        y += Inches(0.6)

    # Bottom value box
    shape = add_rect(slide, Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.6),
                     GOLD_LIGHT, GOLD, Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "传统：高级顾问带人逐格讲，3-5天"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.name = FONT_CN
    add_paragraph(tf, "图谱：初级顾问自己看，3小时", 18, NAVY, bold=True, space_before=6)
    add_paragraph(tf, "→ 释放高级顾问产能，缩短 ramp-up 时间", 18, NAVY, space_before=6)

    add_slide_number(slide, 6)


# ======================== Slide 7: Scenario 2 ========================
def build_slide_7():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1.0), Inches(0.5), Inches(11.333), Inches(0.8),
                "场景 2：方案汇报 — 让客户“看得见”你的专业", 32, NAVY, bold=True)
    add_rect_straight(slide, Inches(1.0), Inches(1.2), Inches(2.0), Pt(3), GOLD)

    # Left - 业务场景
    add_rect_straight(slide, Inches(1.0), Inches(1.6), Inches(5.5), Inches(0.5), NAVY)
    add_textbox(slide, Inches(1.1), Inches(1.62), Inches(5.3), Inches(0.45),
                "业务场景", 18, WHITE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(2.3), Inches(5.2), Inches(2.0),
                "改了 20 个参数，客户问“为什么改这些”。\n传统方式是口头解释或列清单。",
                18, DARK_TEXT)

    # Right - 演示要点
    add_rect_straight(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.5), GOLD)
    add_textbox(slide, Inches(6.9), Inches(1.62), Inches(5.3), Inches(0.45),
                "演示要点", 18, NAVY, bold=True)

    points = [
        "快照对比：基准方案 vs 优化方案",
        "热力图展示变化分布",
        "传播路径可视化：改了 A → 影响 B → 传导到 C",
    ]
    y = Inches(2.3)
    for p in points:
        add_textbox(slide, Inches(7.0), y, Inches(5.0), Inches(0.5),
                    f"•  {p}", 18, DARK_TEXT)
        y += Inches(0.6)

    # Bottom value box
    shape = add_rect(slide, Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.6),
                     GOLD_LIGHT, GOLD, Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "传统：口头解释，客户半信半疑"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.name = FONT_CN
    add_paragraph(tf, "图谱：可视化传播链，客户看得懂", 18, NAVY, bold=True, space_before=6)
    add_paragraph(tf, "→ 汇报更直观，客户信任度更高", 18, NAVY, space_before=6)

    add_slide_number(slide, 7)


# ======================== Slide 8: Scenario 3 ========================
def build_slide_8():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1.0), Inches(0.5), Inches(11.333), Inches(0.8),
                "场景 3：参数调整 — 改了敢交", 32, NAVY, bold=True)
    add_rect_straight(slide, Inches(1.0), Inches(1.2), Inches(2.0), Pt(3), GOLD)

    # Left
    add_rect_straight(slide, Inches(1.0), Inches(1.6), Inches(5.5), Inches(0.5), NAVY)
    add_textbox(slide, Inches(1.1), Inches(1.62), Inches(5.3), Inches(0.45),
                "业务场景", 18, WHITE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(2.3), Inches(5.2), Inches(2.0),
                "客户问“融资成本降 50bp，对整体有什么影响”。\n改完后敢不敢直接交？",
                18, DARK_TEXT)

    # Right
    add_rect_straight(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.5), GOLD)
    add_textbox(slide, Inches(6.9), Inches(1.62), Inches(5.3), Inches(0.45),
                "演示要点", 18, NAVY, bold=True)

    points = [
        "选择场景 → 修改参数 → 重算 → 看结果",
        "左编辑，右结果：改了哪些、影响哪些、影响多大",
        "多个方案并行对比",
    ]
    y = Inches(2.3)
    for p in points:
        add_textbox(slide, Inches(7.0), y, Inches(5.0), Inches(0.5),
                    f"•  {p}", 18, DARK_TEXT)
        y += Inches(0.6)

    # Bottom value box
    shape = add_rect(slide, Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.6),
                     GOLD_LIGHT, GOLD, Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "传统：改完手工抽查，心里没底"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.name = FONT_CN
    add_paragraph(tf, "图谱：自动识别全部影响单元格，100%覆盖", 18, NAVY, bold=True, space_before=6)
    add_paragraph(tf, "→ 交付质量可验证，降低返工风险", 18, NAVY, space_before=6)

    add_slide_number(slide, 8)


# ======================== Slide 9: Scenario 4 ========================
def build_slide_9():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1.0), Inches(0.5), Inches(11.333), Inches(0.8),
                "场景 4：客户问答 — 从翻表到秒答", 32, NAVY, bold=True)
    add_rect_straight(slide, Inches(1.0), Inches(1.2), Inches(2.0), Pt(3), GOLD)

    # Left
    add_rect_straight(slide, Inches(1.0), Inches(1.6), Inches(5.5), Inches(0.5), NAVY)
    add_textbox(slide, Inches(1.1), Inches(1.62), Inches(5.3), Inches(0.45),
                "业务场景", 18, WHITE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(2.3), Inches(5.2), Inches(2.0),
                "客户在群里问“2026年折旧摊销总额是多少”。\n传统：打开 Excel → 找表 → 找行 → 找列 → 回复。",
                18, DARK_TEXT)

    # Right
    add_rect_straight(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.5), GOLD)
    add_textbox(slide, Inches(6.9), Inches(1.62), Inches(5.3), Inches(0.45),
                "演示要点", 18, NAVY, bold=True)

    points = [
        "输入问题 → 秒出结构化回答",
        "指标卡 + 折线图 + 数据表 + 数据来源",
        "相关问题推荐",
    ]
    y = Inches(2.3)
    for p in points:
        add_textbox(slide, Inches(7.0), y, Inches(5.0), Inches(0.5),
                    f"•  {p}", 18, DARK_TEXT)
        y += Inches(0.6)

    # Bottom value box
    shape = add_rect(slide, Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.6),
                     GOLD_LIGHT, GOLD, Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "传统：2-5 分钟找数"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.name = FONT_CN
    add_paragraph(tf, "图谱：秒级响应，带图表", 18, NAVY, bold=True, space_before=6)
    add_paragraph(tf, "→ 客户感知专业度提升，响应速度是竞争力", 18, NAVY, space_before=6)

    add_slide_number(slide, 9)


# ======================== Slide 10: Competitive Barrier ========================
def build_slide_10():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1.0), Inches(0.5), Inches(11.333), Inches(0.8),
                "这不是工具，是你们的竞争壁垒", 32, NAVY, bold=True)
    add_rect_straight(slide, Inches(1.0), Inches(1.2), Inches(2.0), Pt(3), GOLD)

    # Table
    headers = ["维度", "传统咨询公司", "用图谱的咨询公司"]
    rows = [
        ["交付速度", "新人 3-5 天理解模型", "3 小时上手"],
        ["交付质量", "改完手工抽查", "自动全量验证"],
        ["客户体验", "口头/文字解释", "可视化传播链"],
        ["知识沉淀", "文件+人脑", "可复用知识图谱"],
    ]

    col_widths = [Inches(2.5), Inches(4.5), Inches(4.5)]
    col_starts = [Inches(1.0), Inches(3.5), Inches(8.0)]
    row_h = Inches(0.65)
    y = Inches(1.6)

    # Header row
    for j, h in enumerate(headers):
        add_rect_straight(slide, col_starts[j], y, col_widths[j], row_h, NAVY)
        add_textbox(slide, col_starts[j] + Inches(0.1), y + Inches(0.08),
                    col_widths[j] - Inches(0.2), Inches(0.5),
                    h, 18, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    y += row_h

    # Data rows
    for i, row in enumerate(rows):
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        for j, cell in enumerate(row):
            add_rect_straight(slide, col_starts[j], y, col_widths[j], row_h, bg, GRAY, Pt(0.5))
            add_textbox(slide, col_starts[j] + Inches(0.1), y + Inches(0.08),
                        col_widths[j] - Inches(0.2), Inches(0.5),
                        cell, 17, DARK_TEXT, alignment=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
        y += row_h

    # Bottom summary
    add_textbox(slide, Inches(1.0), Inches(4.8), Inches(11.333), Inches(0.5),
                "传统模式：每个项目从零开始 → 成本随规模线性增长",
                20, GRAY)
    add_textbox(slide, Inches(1.0), Inches(5.4), Inches(11.333), Inches(0.5),
                "图谱模式：知识持续积累 → 边际成本递减，壁垒递增",
                20, GOLD, bold=True)

    add_slide_number(slide, 10)


# ======================== Slide 11: Technical Reliability ========================
def build_slide_11():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1.0), Inches(0.5), Inches(11.333), Inches(0.8),
                "技术可靠性", 36, NAVY, bold=True)
    add_rect_straight(slide, Inches(1.0), Inches(1.35), Inches(2.0), Pt(3), GOLD)

    metrics = [
        ("73%", "公式覆盖率，依赖方向100%正确"),
        ("0", "#VALUE! 从 5557 降至 0"),
        ("1e-9", "浮点精度容差机制"),
        ("本地", "本地部署，数据不外传"),
    ]
    card_w = Inches(2.7)
    card_h = Inches(2.2)
    gap = Inches(0.4)
    start_x = Inches(1.0)
    y = Inches(1.9)

    for i, (num, label) in enumerate(metrics):
        x = start_x + i * (card_w + gap)
        # Card
        add_rect(slide, x, y, card_w, card_h, WHITE, GOLD, Pt(2))
        # Number
        add_textbox(slide, x + Inches(0.2), y + Inches(0.3),
                    card_w - Inches(0.4), Inches(0.8),
                    num, 48, NAVY, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_NUM)
        # Label
        add_textbox(slide, x + Inches(0.2), y + Inches(1.2),
                    card_w - Inches(0.4), Inches(0.8),
                    label, 16, DARK_TEXT, alignment=PP_ALIGN.CENTER)

    # Small text below
    add_textbox(slide, Inches(1.0), Inches(4.8), Inches(11.333), Inches(0.4),
                "任何 .xlsx 文件均可解析，不需要模板化", 16, GRAY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 11)


# ======================== Slide 12: Next Steps ========================
def build_slide_12():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1.0), Inches(0.5), Inches(11.333), Inches(0.8),
                "下一步：从小处开始", 36, NAVY, bold=True)
    add_rect_straight(slide, Inches(1.0), Inches(1.35), Inches(2.0), Pt(3), GOLD)

    steps = [
        ("1", "POC 验证（2周）", "拿一个正在做的项目文件，免费解析，看效果"),
        ("2", "试点项目（1-2月）", "选 1-2 个顾问试用，对比效率/质量差异"),
        ("3", "规模化部署", "根据试点效果决定推广范围"),
    ]
    card_w = Inches(3.5)
    card_h = Inches(2.5)
    gap = Inches(0.5)
    start_x = Inches(1.0)
    y = Inches(1.8)

    for i, (num, title, desc) in enumerate(steps):
        x = start_x + i * (card_w + gap)
        # Card
        add_rect(slide, x, y, card_w, card_h, WHITE, GOLD, Pt(2))
        # Number circle
        add_icon_circle(slide, x + Inches(1.35), y + Inches(0.2), Inches(0.6), num)
        # Title
        add_textbox(slide, x + Inches(0.2), y + Inches(0.9),
                    card_w - Inches(0.4), Inches(0.5),
                    title, 20, NAVY, bold=True, alignment=PP_ALIGN.CENTER)
        # Desc
        add_textbox(slide, x + Inches(0.2), y + Inches(1.5),
                    card_w - Inches(0.4), Inches(0.8),
                    desc, 16, DARK_TEXT, alignment=PP_ALIGN.CENTER)

    # Callout
    add_gold_callout(slide, Inches(1.0), Inches(5.0), Inches(11.333), Inches(1.2),
                     "“不需要先做决定。拿一个你们现在的模型文件，我来解析，你们自己看效果。”",
                     20)

    add_slide_number(slide, 12)


# ======================== Slide 13: Thank You ========================
def build_slide_13():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Gold horizontal line
    add_rect_straight(slide, Inches(4.0), Inches(2.2), Inches(5.333), Pt(3), GOLD)

    # Title
    add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10.333), Inches(1.2),
                "感谢聆听", 48, GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    # Gold line
    add_rect_straight(slide, Inches(4.0), Inches(3.7), Inches(5.333), Pt(3), GOLD)

    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10.333), Inches(0.8),
                "期待与您合作", 24, WHITE, alignment=PP_ALIGN.CENTER)

    # Contact placeholder
    add_textbox(slide, Inches(1.5), Inches(6.0), Inches(10.333), Inches(0.5),
                "联系人  /  电话  /  邮箱", 16, GRAY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 13)


# ======================== Build All Slides ========================
build_slide_1()
build_slide_2()
build_slide_3()
build_slide_4()
build_slide_5()
build_slide_6()
build_slide_7()
build_slide_8()
build_slide_9()
build_slide_10()
build_slide_11()
build_slide_12()
build_slide_13()

# Save
OUTPUT = r"D:\project_room\workspace2024\mytest\financial2_neo4j5_claude\财务模型知识图谱_演示V3.pptx"
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"File size: {os.path.getsize(OUTPUT) / 1024:.1f} KB")
