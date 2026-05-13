#!/usr/bin/env python
"""Create a professional PowerPoint presentation for 财务模型知识图谱系统."""

from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Constants ---
NAVY = "1B2A4A"
GOLD = "D4A843"
WHITE = "FFFFFF"
LIGHT_GRAY = "F5F5F5"
DARK_NAVY = "0F1B33"
ACCENT_BLUE = "2E4A7A"
SOFT_GOLD = "F0DCA0"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def rgb(hex_str):
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


NAVY_C = rgb(NAVY)
GOLD_C = rgb(GOLD)
WHITE_C = rgb(WHITE)
LIGHT_GRAY_C = rgb(LIGHT_GRAY)
DARK_NAVY_C = rgb(DARK_NAVY)
ACCENT_BLUE_C = rgb(ACCENT_BLUE)
SOFT_GOLD_C = rgb(SOFT_GOLD)


def add_bg(slide, color_hex):
    """Set solid background color on slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color_hex)


def add_rect(slide, left, top, width, height, fill_hex, line_hex=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shp.line.color.rgb = rgb(line_hex)
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def add_rounded_rect(slide, left, top, width, height, fill_hex, line_hex=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shp.line.color.rgb = rgb(line_hex)
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    return shp


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE_C, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def set_tf_text(tf, text, font_size=18, color=WHITE_C, bold=False,
                font_name="Microsoft YaHei", alignment=PP_ALIGN.LEFT, space_before=0, space_after=0):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return p


def add_slide_number(slide, num, total=12):
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.9), Inches(0.4),
                f"{num}/{total}", font_size=10, color=rgb("8899AA"),
                alignment=PP_ALIGN.RIGHT)


# =====================================================================
# SLIDE 1 — Title
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_NAVY)

# Gold accent line at top
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), GOLD)

# Center title block
add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.2),
            "财务模型知识图谱系统", font_size=48, color=GOLD_C,
            bold=True, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

# Subtitle
add_textbox(slide, Inches(2), Inches(3.3), Inches(9.3), Inches(0.8),
            "给 Excel 财务模型装上 X 光片", font_size=28, color=WHITE_C,
            bold=False, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

# Divider line
add_rect(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.03), GOLD)

# Date
today_str = date.today().strftime("%Y年%m月%d日")
add_textbox(slide, Inches(2), Inches(4.7), Inches(9.3), Inches(0.6),
            today_str, font_size=18, color=rgb("AABBCC"),
            alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

# Bottom gold line
add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), GOLD)

add_slide_number(slide, 1)


# =====================================================================
# SLIDE 2 — Pain Points
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Header bar
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), GOLD)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
            "痛点 — 你们每天面对的问题", font_size=36, color=WHITE_C,
            bold=True, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

pain_points = [
    ("黑盒模型", "打开 50+ Sheet 的财务模型，\n不知道某个数字怎么来的"),
    ("修改恐惧", "改一个参数，不知道会影响\n多少个 downstream 单元格"),
    ("知识不可复用", "模型做完就死在文件里，\n下一个项目从零开始"),
    ("问答靠翻阅", "客户问“折旧摊销总额是多少”，\n需要手动翻表找答案"),
    ("版本对比困难", "两个版本的差异只能看文件 diff，\n不知道影响链"),
]

card_w = Inches(2.2)
card_h = Inches(4.2)
gap = Inches(0.3)
start_x = Inches(0.7)
start_y = Inches(1.5)

for i, (title, desc) in enumerate(pain_points):
    x = start_x + i * (card_w + gap)
    # Card background
    card = add_rounded_rect(slide, x, start_y, card_w, card_h, NAVY, GOLD)
    card.text_frame.paragraphs[0].space_after = Pt(0)

    # Number circle
    num_shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      x + Inches(0.8), start_y + Inches(0.3),
                                      Inches(0.6), Inches(0.6))
    num_shp.fill.solid()
    num_shp.fill.fore_color.rgb = GOLD_C
    num_shp.line.fill.background()
    tf = num_shp.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY_C
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False

    # Title
    add_textbox(slide, x + Inches(0.15), start_y + Inches(1.1), card_w - Inches(0.3), Inches(0.5),
                title, font_size=22, color=GOLD_C, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

    # Description
    add_textbox(slide, x + Inches(0.15), start_y + Inches(1.8), card_w - Inches(0.3), Inches(2.2),
                desc, font_size=16, color=WHITE_C,
                alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

add_slide_number(slide, 2)


# =====================================================================
# SLIDE 3 — Core Concept (Pyramid)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), GOLD)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
            "核心概念 — 知识图谱是什么", font_size=36, color=WHITE_C,
            bold=True, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

# Pyramid layers (top to bottom)
layers = [
    ("Indicator 层", "有业务意义的财务条目", "如“动态总投资”、“营业收入”", 2.5, ACCENT_BLUE),
    ("Table 层", "逻辑表", "如“资金筹措及还本付息表”", 5.0, rgb("3A5F8F")),
    ("Cell 层", "每个单元格，带公式依赖关系", "35,000+ 单元格，73% 公式覆盖率", 7.5, rgb("4A7FAF")),
]

for idx, (name, desc, detail, w, clr) in enumerate(layers):
    y = Inches(1.7) + idx * Inches(1.7)
    x = Inches(6.665 - w / 2)
    shp = add_rounded_rect(slide, x, y, Inches(w), Inches(1.4), str(clr), GOLD)
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{name}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE_C
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(16)
    p2.font.color.rgb = SOFT_GOLD_C
    p2.font.name = "Microsoft YaHei"
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = detail
    p3.font.size = Pt(13)
    p3.font.color.rgb = rgb("CCDDEE")
    p3.font.name = "Microsoft YaHei"
    p3.alignment = PP_ALIGN.CENTER

# Side annotation
add_textbox(slide, Inches(0.5), Inches(2.5), Inches(3.5), Inches(4.0),
            "三层结构\n\n"
            "Cell 层 → 最细粒度\n"
            "每个单元格包含公式、\n"
            "依赖关系、数据值\n\n"
            "Indicator 层 → 业务语义\n"
            "将零散单元格聚合为\n"
            "有业务含义的条目\n\n"
            "Table 层 → 逻辑组织\n"
            "按财务表组织指标，\n"
            "保持业务上下文",
            font_size=15, color=rgb("445566"),
            alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

add_slide_number(slide, 3)


# =====================================================================
# SLIDE 4 — Capability 1: Parse
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), GOLD)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
            "能力 1 — 一键解析，看见模型骨架", font_size=36, color=WHITE_C,
            bold=True, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

metrics = [
    ("49", "张表"),
    ("2,968", "个指标"),
    ("35,000+", "单元格"),
    ("73%", "公式覆盖率"),
]

card_w2 = Inches(2.5)
start_x2 = Inches(1.0)

for i, (num, label) in enumerate(metrics):
    x = start_x2 + i * (card_w2 + Inches(0.35))
    # Card
    card = add_rounded_rect(slide, x, Inches(1.8), card_w2, Inches(3.2), NAVY, GOLD)
    # Number
    add_textbox(slide, x, Inches(2.2), card_w2, Inches(1.5),
                num, font_size=56, color=GOLD_C, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")
    # Label
    add_textbox(slide, x, Inches(3.7), card_w2, Inches(0.6),
                label, font_size=20, color=WHITE_C,
                alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

# Value proposition
vp = add_rounded_rect(slide, Inches(2.5), Inches(5.5), Inches(8.3), Inches(1.2), GOLD, NAVY)
tf = vp.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "价值：10 分钟的工作 → 10 秒完成"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = NAVY_C
p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER

add_slide_number(slide, 4)


# =====================================================================
# SLIDE 5 — Capability 2: Explorer
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), GOLD)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
            "能力 2 — 交互式图谱，追踪依赖链", font_size=36, color=WHITE_C,
            bold=True, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

# Left: description panel
desc_panel = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.2), NAVY, GOLD)
tf = desc_panel.text_frame
tf.word_wrap = True
set_tf_text(tf, "ECharts 交互图谱", 26, GOLD_C, True, "Microsoft YaHei", PP_ALIGN.LEFT, 12, 6)
set_tf_text(tf, "• Sheet → Table → Indicator → Cell 逐层下钻", 18, WHITE_C, False, "Microsoft YaHei", PP_ALIGN.LEFT, 10, 4)
set_tf_text(tf, "• 5 种布局：force / circular / radial / tree / mindmap", 18, WHITE_C, False, "Microsoft YaHei", PP_ALIGN.LEFT, 6, 4)
set_tf_text(tf, "• 自动冻结大规模节点", 18, WHITE_C, False, "Microsoft YaHei", PP_ALIGN.LEFT, 6, 4)
set_tf_text(tf, "• 搜索直达 + 面包屑导航", 18, WHITE_C, False, "Microsoft YaHei", PP_ALIGN.LEFT, 6, 4)
set_tf_text(tf, "• 关系线展示公式依赖方向", 18, WHITE_C, False, "Microsoft YaHei", PP_ALIGN.LEFT, 6, 4)

# Value box
vp2 = add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(4.9), Inches(0.9), GOLD, NAVY)
tf2 = vp2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "回答“这个数字怎么来的”，不再逐格翻阅"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = NAVY_C
p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER

# Right: mock diagram showing hierarchy
hierarchy_items = [
    ("Sheet", "利润表", Inches(6.8), Inches(1.8), Inches(2.2)),
    ("Table", "营业收入明细", Inches(6.8), Inches(3.0), Inches(2.2)),
    ("Indicator", "营业收入", Inches(6.8), Inches(4.2), Inches(2.2)),
    ("Cell", "C15 = SUM(C16:C20)", Inches(6.8), Inches(5.4), Inches(2.2)),
]

# Vertical connector line
add_rect(slide, Inches(7.85), Inches(2.4), Inches(0.04), Inches(3.6), GOLD)

for label, detail, x, y, w in hierarchy_items:
    box = add_rounded_rect(slide, x, y, w, Inches(0.8), NAVY, GOLD)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GOLD_C
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = detail
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE_C
    p2.font.name = "Microsoft YaHei"
    p2.alignment = PP_ALIGN.CENTER

# Arrow indicators between boxes
for i in range(3):
    ay = Inches(2.6) + i * Inches(1.2)
    add_textbox(slide, Inches(7.55), ay, Inches(0.6), Inches(0.4),
                "↓", font_size=20, color=GOLD_C,
                alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

add_slide_number(slide, 5)


# =====================================================================
# SLIDE 6 — Capability 3: Recalc
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), GOLD)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
            "能力 3 — 改参数，看影响", font_size=36, color=WHITE_C,
            bold=True, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

# Left panel: Edit
left_panel = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2), NAVY, GOLD)
tf = left_panel.text_frame
tf.word_wrap = True
set_tf_text(tf, "编辑区", 22, GOLD_C, True, "Microsoft YaHei", PP_ALIGN.LEFT, 10, 8)

# Scenario tabs
scenarios = ["基准", "乐观", "悲观"]
for i, s in enumerate(scenarios):
    tab = add_rounded_rect(slide, Inches(0.8) + i * Inches(1.3), Inches(2.3), Inches(1.1), Inches(0.45),
                           ACCENT_BLUE if i > 0 else GOLD, NAVY)
    tf2 = tab.text_frame
    tf2.paragraphs[0].text = s
    tf2.paragraphs[0].font.size = Pt(14)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = NAVY_C if i == 0 else WHITE_C
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

# Edit items
edit_items = [
    "营业收入增长率: 10% → 12%",
    "折现率: 8% → 7.5%",
    "所得税率: 25% → 15%",
]
for i, item in enumerate(edit_items):
    add_textbox(slide, Inches(0.8), Inches(3.0) + i * Inches(0.6), Inches(5.2), Inches(0.5),
                item, font_size=16, color=WHITE_C,
                alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

# Right panel: Impact
right_panel = add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.2), ACCENT_BLUE, GOLD)
tf = right_panel.text_frame
tf.word_wrap = True
set_tf_text(tf, "影响分析", 22, GOLD_C, True, "Microsoft YaHei", PP_ALIGN.LEFT, 10, 8)

impact_items = [
    ("受影响的指标", "23 个"),
    ("受影响的单元格", "156 个"),
    ("最大变化幅度", "+18.3%"),
    ("NPV 变化", "+12.7%"),
]
for i, (label, val) in enumerate(impact_items):
    y = Inches(2.5) + i * Inches(0.9)
    add_textbox(slide, Inches(7.1), y, Inches(2.5), Inches(0.5),
                label, font_size=16, color=WHITE_C,
                alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")
    add_textbox(slide, Inches(10.0), y, Inches(2.5), Inches(0.5),
                val, font_size=22, color=GOLD_C, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

# Value box
vp3 = add_rounded_rect(slide, Inches(0.8), Inches(5.6), Inches(5.2), Inches(0.8), GOLD, NAVY)
tf2 = vp3.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "改一个数，立刻知道影响范围"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = NAVY_C
p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER

add_slide_number(slide, 6)


# =====================================================================
# SLIDE 7 — Capability 4: Snapshot Compare
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), GOLD)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
            "能力 4 — 快照对比，看清变化传播", font_size=36, color=WHITE_C,
            bold=True, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

# Left: features
left = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.2), NAVY, GOLD)
tf = left.text_frame
tf.word_wrap = True
set_tf_text(tf, "功能特性", 22, GOLD_C, True, "Microsoft YaHei", PP_ALIGN.LEFT, 10, 8)

features = [
    "热力图展示变化分布",
    "传播路径可视化",
    "支持导出对比报告",
    "相对+绝对容差自动滤除浮点噪音",
]
for i, f in enumerate(features):
    add_textbox(slide, Inches(0.8), Inches(2.3) + i * Inches(0.7), Inches(4.9), Inches(0.6),
                f"• {f}", font_size=17, color=WHITE_C,
                alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

# Right: mock heatmap
right = add_rounded_rect(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(5.2), NAVY, GOLD)
tf = right.text_frame
tf.word_wrap = True
set_tf_text(tf, "变化热力图", 22, GOLD_C, True, "Microsoft YaHei", PP_ALIGN.CENTER, 10, 8)

# Mock heatmap grid
import random
random.seed(42)
colors_heat = [rgb("1B2A4A"), rgb("2E4A7A"), rgb("4A7FAF"),
               rgb("D4A843"), rgb("E8553A"), rgb("C0392B")]
for row in range(5):
    for col in range(8):
        v = random.random()
        ci = min(int(v * 6), 5)
        x = Inches(6.8) + col * Inches(0.72)
        y = Inches(2.5) + row * Inches(0.7)
        add_rect(slide, x, y, Inches(0.65), Inches(0.6), str(colors_heat[ci]))

# Legend
legend_items = ["无变化", "微变", "小变", "中变", "大变", "重大"]
for i, lbl in enumerate(legend_items):
    x = Inches(6.8) + i * Inches(0.95)
    y = Inches(6.1)
    add_rect(slide, x, y, Inches(0.4), Inches(0.25), str(colors_heat[i]))
    add_textbox(slide, x, y + Inches(0.28), Inches(0.9), Inches(0.3),
                lbl, font_size=10, color=rgb("AABBCC"),
                alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

# Value box
vp4 = add_rounded_rect(slide, Inches(0.8), Inches(5.6), Inches(5.0), Inches(0.8), GOLD, NAVY)
tf2 = vp4.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "版本对比不再是“肉眼找不同”"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = NAVY_C
p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER

add_slide_number(slide, 7)


# =====================================================================
# SLIDE 8 — LLM Q&A
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), GOLD)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
            "进阶能力 — LLM 智能问答", font_size=36, color=WHITE_C,
            bold=True, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

# Left: features
left = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.2), NAVY, GOLD)
tf = left.text_frame
tf.word_wrap = True
set_tf_text(tf, "核心能力", 22, GOLD_C, True, "Microsoft YaHei", PP_ALIGN.LEFT, 10, 8)

qa_features = [
    "置信度评分 — 每个回答附带可靠度指标",
    "指标卡 + 折线图 + 数据表 — 结构化展示",
    "数据来源卡片可追溯 — 知其然，知其所以然",
    "中文分词 + 模糊匹配 — 支持自然语言提问",
    "流式输出 — 实时生成回答内容",
]
for i, f in enumerate(qa_features):
    add_textbox(slide, Inches(0.8), Inches(2.3) + i * Inches(0.65), Inches(4.9), Inches(0.6),
                f"• {f}", font_size=16, color=WHITE_C,
                alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

# Right: sample Q&A
right = add_rounded_rect(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(5.2), ACCENT_BLUE, GOLD)
tf = right.text_frame
tf.word_wrap = True
set_tf_text(tf, "示例问答", 22, GOLD_C, True, "Microsoft YaHei", PP_ALIGN.LEFT, 10, 8)

# Question 1
q1 = add_rounded_rect(slide, Inches(6.8), Inches(2.4), Inches(5.7), Inches(1.0), NAVY, GOLD)
tf = q1.text_frame
tf.word_wrap = True
set_tf_text(tf, "问：", 14, GOLD_C, True, "Microsoft YaHei", PP_ALIGN.LEFT, 6, 0)
set_tf_text(tf, "2026年的营业收入是多少？", 16, WHITE_C, False, "Microsoft YaHei", PP_ALIGN.LEFT, 2, 4)

# Answer 1 mock
a1 = add_rounded_rect(slide, Inches(6.8), Inches(3.6), Inches(5.7), Inches(1.3), DARK_NAVY, GOLD)
tf = a1.text_frame
tf.word_wrap = True
set_tf_text(tf, "答：2026年营业收入为 12,500 万元", 15, SOFT_GOLD_C, True, "Microsoft YaHei", PP_ALIGN.LEFT, 6, 2)
set_tf_text(tf, "置信度：95%  |  数据来源：利润表.C15", 12, rgb("88AACC"), False, "Microsoft YaHei", PP_ALIGN.LEFT, 2, 2)
set_tf_text(tf, "同比增长：+12.3%", 12, rgb("4CAF50"), False, "Microsoft YaHei", PP_ALIGN.LEFT, 2, 0)

# Question 2
q2 = add_rounded_rect(slide, Inches(6.8), Inches(5.1), Inches(5.7), Inches(1.0), NAVY, GOLD)
tf = q2.text_frame
tf.word_wrap = True
set_tf_text(tf, "问：", 14, GOLD_C, True, "Microsoft YaHei", PP_ALIGN.LEFT, 6, 0)
set_tf_text(tf, "折旧摊销的费用构成是什么？", 16, WHITE_C, False, "Microsoft YaHei", PP_ALIGN.LEFT, 2, 4)

add_slide_number(slide, 8)


# =====================================================================
# SLIDE 9 — Architecture
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), GOLD)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
            "技术架构", font_size=36, color=WHITE_C,
            bold=True, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

# Data flow boxes
flow_items = [
    ("Excel 文件", NAVY),
    ("解析引擎", ACCENT_BLUE),
    ("知识图谱", rgb("3A5F8F")),
    ("重算引擎", rgb("4A7FAF")),
    ("应用层", GOLD),
]

box_w = Inches(2.0)
box_h = Inches(1.0)
arrow_w = Inches(0.6)
total_w = 5 * box_w + 4 * arrow_w
start_x3 = (SLIDE_W - total_w) / 2
y_flow = Inches(1.8)

for i, (label, clr) in enumerate(flow_items):
    x = start_x3 + i * (box_w + arrow_w)
    shp = add_rounded_rect(slide, x, y_flow, box_w, box_h, str(clr), GOLD)
    tf = shp.text_frame
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE_C if clr != GOLD else NAVY_C
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Arrow
    if i < 4:
        arrow_x = x + box_w
        add_textbox(slide, arrow_x, y_flow + Inches(0.25), arrow_w, Inches(0.5),
                    "→", font_size=28, color=GOLD_C,
                    alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

# Sub-items below each flow box
sub_items = [
    [".xlsx/.xls", "多 Sheet", "公式解析"],
    ["openpyxl", "formulas库", "依赖图"],
    ["Neo4j 图数据库", "Cell/Indicator/Table", "关系边"],
    ["依赖方向重算", "#VALUE! 5557→0", "浮点容差"],
    ["图谱浏览", "快照对比", "LLM问答"],
]

for i, subs in enumerate(sub_items):
    x = start_x3 + i * (box_w + arrow_w)
    for j, sub in enumerate(subs):
        add_textbox(slide, x, y_flow + Inches(1.2) + j * Inches(0.4), box_w, Inches(0.4),
                    sub, font_size=12, color=rgb("667788"),
                    alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

# Key stats at bottom
stats_panel = add_rounded_rect(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.5), NAVY, GOLD)
tf = stats_panel.text_frame
tf.word_wrap = True
set_tf_text(tf, "关键指标", 18, GOLD_C, True, "Microsoft YaHei", PP_ALIGN.CENTER, 8, 4)
set_tf_text(tf, "公式解析率：73% 公式覆盖  |  依赖方向 100% 正确  |  重算精度：#VALUE! 从 5557 降至 0  |  浮点精度：相对+绝对容差", 14, WHITE_C, False, "Microsoft YaHei", PP_ALIGN.CENTER, 4, 0)

add_slide_number(slide, 9)


# =====================================================================
# SLIDE 10 — Business Value
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), GOLD)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
            "业务价值", font_size=36, color=WHITE_C,
            bold=True, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

# Value table
table_data = [
    ("价值维度", "传统方式", "知识图谱方式", "提升"),
    ("模型理解时间", "数小时", "数分钟", "10x+"),
    ("修改影响分析", "人工追踪", "自动可视化", "精准全面"),
    ("知识资产沉淀", "文件死存", "可复用图谱", "沉淀复用"),
    ("客户响应速度", "翻表找数", "即时问答", "实时响应"),
]

tbl_left = Inches(1.5)
tbl_top = Inches(1.6)
col_widths = [Inches(2.8), Inches(2.8), Inches(2.8), Inches(2.0)]
row_height = Inches(0.85)

table_shape = slide.shapes.add_table(len(table_data), 4, tbl_left, tbl_top, sum(col_widths), row_height * len(table_data))
table = table_shape.table

for ci, w in enumerate(col_widths):
    table.columns[ci].width = w

for ri, row in enumerate(table_data):
    for ci, cell_text in enumerate(row):
        cell = table.cell(ri, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = cell_text
        p.font.size = Pt(16)
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        if ri == 0:  # Header
            p.font.bold = True
            p.font.color.rgb = NAVY_C
            cell.fill.solid()
            cell.fill.fore_color.rgb = GOLD_C
        else:
            p.font.color.rgb = rgb("334455") if ci < 2 else GOLD_C
            if ci >= 2:
                p.font.bold = True
            # Alternate row shading
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb("EEF2F7")
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE_C

# Decorative side element
add_rect(slide, Inches(11.5), Inches(2.0), Inches(0.06), Inches(3.5), GOLD)

add_textbox(slide, Inches(11.8), Inches(2.0), Inches(1.2), Inches(3.5),
            "数据驱动\n决策效率\n提升",
            font_size=20, color=NAVY_C, bold=True,
            alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

add_slide_number(slide, 10)


# =====================================================================
# SLIDE 11 — Cooperation Directions
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), GOLD)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
            "合作方向", font_size=36, color=WHITE_C,
            bold=True, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

coop_items = [
    ("01", "模型标准化",
     "将常用模板转化为可复用的知识图谱",
     "建立行业标准模板库，降低新项目启动成本"),
    ("02", "定制开发",
     "针对特定行业/场景扩展解析和问答能力",
     "基于客户具体需求，定制化部署与集成"),
    ("03", "服务集成",
     "嵌入现有咨询工作流，提升团队效率",
     "与现有工具链接，无缝升级工作方式"),
]

for i, (num, title, desc, detail) in enumerate(coop_items):
    x = Inches(0.8) + i * Inches(4.1)
    y = Inches(1.8)

    # Number
    num_shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.8), Inches(0.8))
    num_shp.fill.solid()
    num_shp.fill.fore_color.rgb = GOLD_C
    num_shp.line.fill.background()
    tf = num_shp.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY_C
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Card
    card = add_rounded_rect(slide, x, y + Inches(1.0), Inches(3.7), Inches(4.2), NAVY, GOLD)
    tf = card.text_frame
    tf.word_wrap = True
    set_tf_text(tf, title, 24, GOLD_C, True, "Microsoft YaHei", PP_ALIGN.LEFT, 10, 8)
    set_tf_text(tf, desc, 16, WHITE_C, False, "Microsoft YaHei", PP_ALIGN.LEFT, 6, 8)
    set_tf_text(tf, detail, 14, rgb("AABBCC"), False, "Microsoft YaHei", PP_ALIGN.LEFT, 6, 0)

add_slide_number(slide, 11)


# =====================================================================
# SLIDE 12 — Thank You
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), GOLD)

add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.2),
            "感谢聆听", font_size=52, color=GOLD_C,
            bold=True, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

add_rect(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.03), GOLD)

add_textbox(slide, Inches(2), Inches(3.8), Inches(9.3), Inches(0.8),
            "期待与您合作", font_size=28, color=WHITE_C,
            bold=False, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

# Contact placeholder
contact = add_rounded_rect(slide, Inches(4.0), Inches(5.0), Inches(5.3), Inches(1.5), ACCENT_BLUE, GOLD)
tf = contact.text_frame
tf.word_wrap = True
set_tf_text(tf, "联系我们", 18, GOLD_C, True, "Microsoft YaHei", PP_ALIGN.CENTER, 10, 4)
set_tf_text(tf, "email@example.com  |  +86 XXX XXXX XXXX", 14, WHITE_C, False, "Microsoft YaHei", PP_ALIGN.CENTER, 4, 0)
set_tf_text(tf, "www.example.com", 14, rgb("AABBCC"), False, "Microsoft YaHei", PP_ALIGN.CENTER, 4, 0)

add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), GOLD)

add_slide_number(slide, 12)


# =====================================================================
# Save
# =====================================================================
output_path = r"D:\project_room\workspace2024\mytest\financial2_neo4j5_claude\财务模型知识图谱_演示.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
